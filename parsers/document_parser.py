# parsers/document_parser.py
from pptx import Presentation
import openpyxl
from transformers import pipeline
import torch

def summarize_text(text: str) -> str:
    """Generates a concise summary of a long block of text using an AI model."""
    try:
        summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")
        truncated_text = text[:1024]
        summary = summarizer(truncated_text, max_length=60, min_length=20, do_sample=False)
        result = summary[0]['summary_text']

        # --- Clean up GPU memory ---
        print("Cleaning up summarizer model from GPU memory...")
        del summarizer
        import gc
        gc.collect()
        torch.cuda.empty_cache()

        return result
    except Exception as e:
        print(f"Summarization failed, falling back to truncation: {e}")
        return (text[:150] + '...') if len(text) > 150 else text

def parse_document(file_path: str) -> dict:
    """Parses a document (.pptx or .xlsx) to extract raw text and an AI-generated description."""
    raw_text = ""
    description = ""
    all_text = []

    try:
        if file_path.endswith('.pptx'):
            print(f"Parsing PowerPoint file: {file_path}")
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            
        elif file_path.endswith('.xlsx'):
            print(f"Parsing Excel file: {file_path}")
            workbook = openpyxl.load_workbook(file_path)
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            all_text.append(str(cell.value))
        else:
            raise ValueError("Unsupported file type for document parser.")

        raw_text = "\n".join(all_text).strip()
        
        if raw_text:
            description = summarize_text(raw_text)
        else:
            description = "Document appears to be empty or contains no text."

        print("Document parsing successful.")
        return {
            "description": description.strip(),
            "raw_text": raw_text
        }

    except Exception as e:
        print(f"An error occurred in document parser: {e}")
        return {
            "description": f"ERROR: Document processing failed. Details: {e}",
            "raw_text": ""
        }
